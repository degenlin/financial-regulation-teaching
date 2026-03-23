
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = "C:/Users/Administrator/Downloads/金融监管理论与实践/PPT课件"

C_DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x6B)
C_MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
C_ACCENT     = RGBColor(0xC0, 0x39, 0x2B)
C_GOLD       = RGBColor(0xD4, 0xAC, 0x0D)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
C_DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)
C_GREEN      = RGBColor(0x1E, 0x8B, 0x4C)
C_PALE_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)
C_SOFT_BLUE  = RGBColor(0x9D, 0xC3, 0xE6)
C_NIGHT      = RGBColor(0x0D, 0x1F, 0x47)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT):
    b = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

# ── slide makers ─────────────────────────────────────────────

def cover(prs, num, title, subtitle, desc):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BLUE)
    rect(sl, 0, 0, 0.25, 7.5, C_GOLD)
    rect(sl, 0, 6.6, 13.33, 0.9, C_MID_BLUE)
    rect(sl, 0.25, 6.6, 13.08, 0.07, C_GOLD)
    txt(sl, "王娴《金融监管理论与实践》配套课件",
        0.6, 0.3, 12.0, 0.5, 13, color=C_GOLD)
    txt(sl, f"第 {num} 章",
        0.6, 1.1, 5, 0.85, 34, bold=True, color=C_WHITE)
    txt(sl, title,
        0.6, 2.0, 12.0, 1.15, 38, bold=True, color=C_GOLD)
    txt(sl, subtitle,
        0.6, 3.3, 12.0, 0.65, 19, color=C_PALE_BLUE)
    txt(sl, desc,
        0.6, 4.1, 11.5, 1.2, 13, color=C_SOFT_BLUE)
    txt(sl, "金融监管理论与实践  |  教学课件",
        0.6, 6.65, 8, 0.45, 11, color=C_WHITE)

def objectives(prs, num, items):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_WHITE)
    rect(sl, 0, 0, 13.33, 1.5, C_DARK_BLUE)
    rect(sl, 0, 1.5, 13.33, 0.06, C_GOLD)
    txt(sl, f"第{num}章", 0.3, 0.1, 3, 0.45, 11,
        color=C_GOLD, bold=True)
    txt(sl, "本章学习目标", 0.3, 0.5, 10, 0.82,
        28, bold=True, color=C_WHITE)
    y = 1.82
    for i, ob in enumerate(items):
        rect(sl, 0.4, y, 0.55, 0.5, C_MID_BLUE)
        txt(sl, str(i+1), 0.4, y, 0.55, 0.5,
            17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, ob, 1.1, y, 11.5, 0.5, 14,
            color=C_DARK_BLUE)
        y += 0.72

def content(prs, num, title, bullets, accent=C_MID_BLUE):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_WHITE)
    rect(sl, 0, 0, 13.33, 1.4, C_DARK_BLUE)
    rect(sl, 0, 1.4, 0.12, 6.1, accent)
    txt(sl, f"第{num}章", 0.2, 0.1, 2, 0.45, 11,
        color=C_GOLD, bold=True)
    txt(sl, title, 0.2, 0.45, 12.5, 0.82,
        24, bold=True, color=C_WHITE)
    y = 1.62
    for item in bullets:
        lv, text = item
        indent = 0.32 + lv * 0.38
        mk = "◆" if lv == 0 else "▸"
        sz = 13 if lv == 0 else 11
        c_txt = C_DARK_BLUE if lv == 0 else C_DARK_GRAY
        txt(sl, mk, indent, y, 0.38, 0.4,
            sz, bold=(lv == 0),
            color=accent if lv == 0 else C_DARK_GRAY)
        txt(sl, text, indent+0.33, y, 12.2-indent, 0.4,
            sz, bold=(lv == 0), color=c_txt)
        y += 0.44 if lv == 0 else 0.39

def case(prs, num, title, content_txt, analysis):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_WHITE)
    rect(sl, 0, 0, 13.33, 1.4, C_ACCENT)
    rect(sl, 0, 1.4, 13.33, 0.06, C_GOLD)
    txt(sl, f"第{num}章  |  案例分析", 0.3, 0.1, 8, 0.42,
        12, color=C_GOLD, bold=True)
    txt(sl, title, 0.3, 0.5, 12.5, 0.78, 24,
        bold=True, color=C_WHITE)
    rect(sl, 0.3, 1.62, 12.7, 2.28, C_LIGHT_GRAY)
    rect(sl, 0.3, 1.62, 0.08, 2.28, C_ACCENT)
    txt(sl, "【案例背景】", 0.55, 1.66, 3, 0.38,
        13, bold=True, color=C_ACCENT)
    txt(sl, content_txt, 0.55, 2.05, 12.3, 1.78,
        12, color=C_DARK_GRAY)
    txt(sl, "【案例分析】", 0.3, 4.06, 3, 0.38,
        13, bold=True, color=C_DARK_BLUE)
    y = 4.5
    for a in analysis:
        txt(sl, f"▸  {a}", 0.4, y, 12.5, 0.4, 12,
            color=C_DARK_BLUE)
        y += 0.42

def summary(prs, num, title, points):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_WHITE)
    rect(sl, 0, 0, 13.33, 1.4, C_GREEN)
    rect(sl, 0, 1.4, 13.33, 0.06, C_GOLD)
    txt(sl, f"第{num}章  |  本章小结", 0.3, 0.1, 8, 0.42,
        12, color=C_GOLD, bold=True)
    txt(sl, f"{title} — 核心要点", 0.3, 0.5, 12, 0.78,
        24, bold=True, color=C_WHITE)
    y = 1.62
    for i, pt in enumerate(points):
        bg = C_LIGHT_GRAY if i % 2 == 0 else C_WHITE
        rect(sl, 0.3, y, 12.7, 0.6, bg)
        rect(sl, 0.3, y, 0.06, 0.6, C_GREEN)
        txt(sl, f"  {pt}", 0.5, y+0.08, 12.3, 0.44,
            13, color=C_DARK_BLUE)
        y += 0.63

def questions(prs, num, qs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_DARK_BLUE)
    rect(sl, 0, 0, 13.33, 1.5, C_NIGHT)
    rect(sl, 0, 1.5, 13.33, 0.06, C_GOLD)
    txt(sl, f"第{num}章  |  思考与讨论", 0.3, 0.1, 8, 0.42,
        12, color=C_GOLD, bold=True)
    txt(sl, "本章思考题", 0.3, 0.5, 10, 0.82,
        28, bold=True, color=C_WHITE)
    y = 1.78
    for i, q in enumerate(qs):
        rect(sl, 0.3, y, 0.68, 0.52, C_GOLD)
        txt(sl, f"Q{i+1}", 0.3, y, 0.68, 0.52,
            16, bold=True, color=C_DARK_BLUE,
            align=PP_ALIGN.CENTER)
        txt(sl, q, 1.13, y, 11.8, 0.52, 13,
            color=C_PALE_BLUE)
        y += 0.76

# ── main ─────────────────────────────────────────────────────

def make_ppt(ch):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    n = ch["num"]
    cover(prs, n, ch["title"], ch["subtitle"], ch["description"])

    sec_titles = [s["title"] for s in ch["sections"]]
    content(prs, n, "本章内容框架",
            [[0, t] for t in sec_titles], accent=C_GOLD)

    objectives(prs, n, ch["objectives"])

    for sec in ch["sections"]:
        content(prs, n, sec["title"], sec["bullets"])

    case(prs, n, ch["case_title"],
         ch["case_content"], ch["case_analysis"])

    summary(prs, n, ch["title"], ch["summary"])

    questions(prs, n, ch["questions"])

    out = os.path.join(OUTPUT_DIR, ch["file"])
    prs.save(out)
    print(f"  OK  {ch['file']}  ({len(prs.slides)} pages)")


if __name__ == "__main__":
    data_path = "C:/Users/Administrator/Downloads/金融监管理论与实践/chapters_data.json"
    with open(data_path, encoding="utf-8") as f:
        data = json.load(f)
    print("=" * 55)
    print("  Generating PPTs ...")
    print("=" * 55)
    for ch in data["chapters"]:
        make_ppt(ch)
    print("\nDone. Output:", OUTPUT_DIR)
